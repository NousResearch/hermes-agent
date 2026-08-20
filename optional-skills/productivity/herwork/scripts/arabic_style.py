"""Arabic typography defaults for HerWork deliverables.

Cairo is the house font for anything that contains Arabic — the office
suite default (Calibri) falls back to an ugly Arabic rendering. These
helpers set Cairo on docx / pptx / pdf output, including the
complex-script font slot that Arabic text is actually shaped from.

For docx and pptx only the font NAME is embedded — the operating system
resolves it, so those work anywhere Cairo (or a fallback) is installed.
PDF generation via reportlab embeds the font file itself, so
``register_pdf_font`` locates the TTF across Linux/macOS/Windows font
directories and tells you to install Cairo if it can't.

Usage:
    from arabic_style import style_docx, style_pptx, register_pdf_font, FONT

    style_docx(doc)                      # after building the Document
    style_pptx(prs)                      # after building the Presentation
    register_pdf_font()                  # then canvas.setFont(FONT, size)
"""
import os
from pathlib import Path

FONT = "Cairo"

_DOCX_STYLES = (
    "Normal", "Title", "Heading 1", "Heading 2", "Heading 3",
    "List Bullet", "List Number",
)


def _font_dirs():
    """Font directories across Linux, macOS, and Windows."""
    home = Path.home()
    dirs = [
        home / ".local/share/fonts",          # Linux (user)
        home / ".fonts",                      # Linux (legacy user)
        Path("/usr/local/share/fonts"),       # Linux (local system)
        Path("/usr/share/fonts"),             # Linux (system)
        home / "Library/Fonts",               # macOS (user)
        Path("/Library/Fonts"),               # macOS (system)
    ]
    local_appdata = os.environ.get("LOCALAPPDATA")
    if local_appdata:                         # Windows 10+ per-user fonts
        dirs.append(Path(local_appdata) / "Microsoft/Windows/Fonts")
    windir = os.environ.get("WINDIR", r"C:\Windows")
    dirs.append(Path(windir) / "Fonts")       # Windows system fonts
    return dirs


def find_font_ttf(family: str = FONT) -> str:
    """Locate the TTF for an installed font family, cross-platform.

    Exact-family files (``Cairo-*.ttf``) win over lookalikes
    (``CairoPlay-*.ttf``). Raises FileNotFoundError with an install hint
    when the family isn't installed."""
    for base in _font_dirs():
        if not base.is_dir():
            continue
        hits = sorted(
            base.rglob(f"{family}*.ttf"),
            key=lambda p: (not p.name.startswith(f"{family}-"), len(p.name)),
        )
        if hits:
            return str(hits[0])
    raise FileNotFoundError(
        f"{family} font not found. Install it first, e.g. from "
        f"https://fonts.google.com/specimen/{family.replace(' ', '+')} "
        f"(or your OS package manager), then retry."
    )


def style_docx(doc) -> None:
    """Apply Cairo to a python-docx Document's base + heading styles.

    Sets the Latin (ascii/hAnsi) AND complex-script (cs) slots — Arabic is
    shaped from the cs slot, so setting only ``font.name`` leaves Arabic on
    the fallback font."""
    from docx.oxml.ns import qn

    for name in _DOCX_STYLES:
        try:
            style = doc.styles[name]
        except KeyError:
            continue
        style.font.name = FONT
        rfonts = style.element.get_or_add_rPr().get_or_add_rFonts()
        for attr in ("w:ascii", "w:hAnsi", "w:cs"):
            rfonts.set(qn(attr), FONT)


def _pptx_run_font(run) -> None:
    from pptx.oxml.ns import qn

    run.font.name = FONT  # sets a:latin
    rPr = run._r.get_or_add_rPr()
    cs = rPr.find(qn("a:cs"))
    if cs is None:
        cs = rPr.makeelement(qn("a:cs"), {})
        rPr.append(cs)
    cs.set("typeface", FONT)


def style_pptx(prs) -> None:
    """Set Cairo on every text run in a python-pptx Presentation,
    including table cells."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        _pptx_run_font(run)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                _pptx_run_font(run)


def register_pdf_font() -> str:
    """Register Cairo with reportlab; returns the font name to setFont()."""
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    pdfmetrics.registerFont(TTFont(FONT, find_font_ttf()))
    return FONT
