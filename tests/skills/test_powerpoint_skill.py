"""Tests for the pptx overflow estimator (issue #91967: pptx_create.py
silently overfills dense body-placeholder content).

Behavior contracts, not snapshots: assertions pin relations that must
hold (warnings fire on known-overflow corpora, stay silent on known-safe
ones, inherited sizes dominate explicit ones), verified against the
properties of decks the create script actually writes.
"""
import json
import os
import subprocess
import sys
from pathlib import Path

import pytest

from pptx import Presentation

REPO = Path(__file__).resolve().parents[2]
SKILL = REPO / "skills" / "productivity" / "powerpoint"
SCRIPTS = SKILL / "scripts"
sys.path.insert(0, str(SCRIPTS))

from pptx_overflow import (  # noqa: E402
    INHERITED_INDENT_IN,
    INHERITED_SIZE_PT,
    effective_size_pt,
    estimate_bullets_overflow,
    paragraph_height_pt,
    text_width_pt,
    wrapped_line_count,
)


def run(script, *args):
    env = dict(os.environ, LC_ALL="C", PYTHONIOENCODING="utf-8")
    proc = subprocess.run(
        [sys.executable, str(SCRIPTS / script), *args],
        capture_output=True, text=True, encoding="utf-8", env=env)
    assert proc.returncode == 0, f"{script} failed: {proc.stderr}"
    return json.loads(proc.stdout)


def write_spec(tmp_path, spec):
    path = tmp_path / "spec.json"
    with open(path, "w", encoding="utf-8") as fh:
        json.dump(spec, fh)
    return str(path)


# --- inherited sizing matches the default-template master -------------------

def test_inherited_sizes_match_default_master():
    # Ground truth from the default python-pptx template's slide master
    # bodyStyle (lvl1pPr..lvl5pPr defRPr sz): 3200/2800/2400/2000/2000.
    prs = Presentation()
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
          "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    body_style = prs.slide_masters[0].element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}txStyles"
    ).find("p:bodyStyle", ns)
    for level, size in INHERITED_SIZE_PT.items():
        lvl = body_style.find(f"a:lvl{level + 1}pPr", ns)
        rpr = lvl.find("a:defRPr", ns)
        actual = int(rpr.get("sz")) / 100.0
        assert actual == size, (
            f"level {level}: estimator assumes {size}pt, master says "
            f"{actual}pt — update INHERITED_SIZE_PT")
        actual_indent = int(lvl.get("marL")) / 914400.0
        assert actual_indent == pytest.approx(
            INHERITED_INDENT_IN[level], abs=0.0001
        ), (
            f"level {level}: estimator assumes "
            f"{INHERITED_INDENT_IN[level]}in, master says "
            f"{actual_indent}in — update INHERITED_INDENT_IN"
        )


def test_title_content_layout_does_not_override_master_body_levels():
    # The estimator uses master bodyStyle sizes/indents only for this
    # layout. Pin the OOXML cascade assumption: its content placeholder's
    # list style must not introduce per-level overrides.
    prs = Presentation()
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
          "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    body = next(ph for ph in prs.slide_layouts[1].placeholders
                if ph.placeholder_format.idx == 1)
    layout_style = body._element.find("p:txBody/a:lstStyle", ns)
    assert layout_style is not None
    assert not any(child.tag.rsplit("}", 1)[-1].startswith("lvl")
                   for child in layout_style)


def test_effective_size_uses_level_inheritance_when_unspecified():
    assert effective_size_pt({"text": "x"}) == 32.0
    assert effective_size_pt({"text": "x", "level": 1}) == 28.0
    assert effective_size_pt("bare string") == 32.0


def test_explicit_size_overrides_inherited():
    assert effective_size_pt({"text": "x", "size": 14}) == 14.0


def test_size_null_and_zero_fall_back_to_inherited():
    # style_run ignores falsy size; the estimator must agree (no TypeError).
    assert effective_size_pt({"text": "x", "size": None}) == 32.0
    assert effective_size_pt({"text": "x", "size": 0}) == 32.0


def test_inherited_size_is_larger_than_old_assumption():
    # Regression guard: level-0 inherited text is NOT 18pt; assuming it was
    # caused false negatives on the motivating case.
    assert effective_size_pt({"text": "x"}) > 18.0


# --- wrapping behavior -------------------------------------------------------

def test_short_bullet_is_single_line():
    assert wrapped_line_count("Overview", 612.0) == 1


def test_long_bullet_wraps_to_multiple_lines():
    text = ("First major point about quarterly infrastructure spending "
            "trends across every region we operate in this fiscal year, "
            "including the carry-over commitments from last year's plan")
    assert wrapped_line_count(text, 612.0) >= 3


def test_hard_line_breaks_force_new_lines():
    assert wrapped_line_count("a\nb\nc", 6000.0) == 3


def test_cjk_text_wraps_by_character_width():
    # CJK glyphs are ~1em wide: 10 chars at 32pt ~ 320pt > 100pt frame,
    # and no spaces means wrap must still happen (char-level fallback).
    assert wrapped_line_count("漢字" * 10, 100.0, 32.0) >= 3


def test_cjk_char_is_full_width():
    assert text_width_pt("漢", 32.0) == 32.0


def test_oversized_unbreakable_token_gets_split():
    assert wrapped_line_count("x" * 200, 100.0) > 1


def test_empty_text_still_one_line():
    assert wrapped_line_count("   ", 612.0) == 1


def test_level_indent_reduces_usable_width():
    # Same explicit size at both levels: deeper indent must wrap to more
    # lines (or equal, never fewer).
    flat = paragraph_height_pt("x " * 60, size_pt=18.0, level=0)
    indented = paragraph_height_pt("x " * 60, size_pt=18.0, level=1)
    assert indented >= flat


# --- estimation decisions ----------------------------------------------------

def test_no_warning_for_typical_safe_slide():
    bullets = ["Overview", "Details", {"text": "Sub", "level": 1}]
    assert estimate_bullets_overflow(bullets) is None


def test_motivating_case_now_warns():
    # The issue reproduction: ~11 mixed bullets at inherited sizes in the
    # default placeholder. Under the correct 32pt inheritance this is
    # borderline-to-over; the estimator must not silently pass it.
    bullets = [
        "First major point about quarterly infrastructure spending trends",
        "Second major point covering headcount growth and onboarding costs",
        "Third point walking through the vendor consolidation plan timeline",
        "Fourth point detailing security review findings and remediation",
        "Fifth point summarizing the budget ask and phased approval gates",
        "Sixth point listing risks, mitigations, and open action owners",
        {"text": "sub-item alpha with additional detail", "level": 1},
        {"text": "sub-item beta with additional detail", "level": 1},
        "Ninth point recapping decisions needed today",
    ]
    est = estimate_bullets_overflow(bullets)
    assert est is not None, "motivating overflow case must warn"


def test_dense_slide_warns_with_heights_ordered():
    bullets = [f"Bullet {i} with a fairly long sentence of operational "
               "detail that wraps onto a second line" for i in range(12)]
    est = estimate_bullets_overflow(bullets)
    assert est is not None
    assert est["estimated_text_height_pt"] > est["frame_usable_height_pt"]
    assert "auto-shrunk or clipped" in est["hint"]


def test_smaller_font_monotonically_reduces_estimate():
    text = f"Bullet {{i}} with plenty of wrapping detail text"
    small = sum(paragraph_height_pt(text.format(i=i), size_pt=12)
                for i in range(10))
    big = sum(paragraph_height_pt(text.format(i=i), size_pt=32)
              for i in range(10))
    assert small < big


def test_none_when_no_bullets():
    assert estimate_bullets_overflow([]) is None


# --- end-to-end through the create script ------------------------------------

def test_create_reports_overflow_warning(tmp_path):
    spec = {
        "slide_size": "16:9",
        "slides": [
            {"layout": "title", "title": "Deck"},
            {"layout": "title_content", "title": "Dense",
             "bullets": [f"Bullet {i}: long operational detail sentence "
                         "that wraps in the placeholder frame and then "
                         "keeps going onto another rendered line"
                         for i in range(16)]},
        ],
    }
    result = run("pptx_create.py", write_spec(tmp_path, spec),
                 str(tmp_path / "out.pptx"))
    assert result["ok"] is True
    warnings = result["overflow_warnings"]  # key always present
    assert warnings, "dense deck must warn"
    first = warnings[0]
    assert first["slide_index"] == 1          # zero-based
    assert first["estimated_text_height_pt"] > first["frame_usable_height_pt"]


def test_create_clean_deck_reports_empty_warnings(tmp_path):
    spec = {
        "slide_size": "16:9",
        "slides": [
            {"layout": "title_content", "title": "Lean",
             "bullets": ["One", "Two", {"text": "Two-a", "level": 1}]},
        ],
    }
    result = run("pptx_create.py", write_spec(tmp_path, spec),
                 str(tmp_path / "clean.pptx"))
    assert result["ok"] is True
    assert result["overflow_warnings"] == []


def test_blank_layout_textbox_fallback_does_not_warn(tmp_path):
    # Blank layout has no body placeholder -> auto-fitting textbox, which
    # grows instead of clipping; the guardrail must stay out of its way.
    spec = {
        "slide_size": "16:9",
        "slides": [
            {"layout": "blank",
             "bullets": [f"Bullet {i} long enough to be flagged if this "
                         "were a fixed placeholder frame" for i in range(20)]},
        ],
    }
    result = run("pptx_create.py", write_spec(tmp_path, spec),
                 str(tmp_path / "blank.pptx"))
    assert result["ok"] is True
    assert result["overflow_warnings"] == []


def test_non_title_content_layouts_are_not_estimated(tmp_path):
    # Only title_content uses the master-body sizing model; other layouts
    # override level styles in their layout XML, so the guardrail must
    # stay silent there rather than emit false positives.
    dense = [{"layout": layout, "title": "Dense",
              "bullets": [f"Bullet {i} would be flagged if estimated "
                          "with master sizes" for i in range(20)]}
             for layout in ("title", "section", "two_content")]
    spec = {"slide_size": "16:9", "slides": dense}
    result = run("pptx_create.py", write_spec(tmp_path, spec),
                 str(tmp_path / "layouts.pptx"))
    assert result["ok"] is True
    assert result["overflow_warnings"] == []


def test_unknown_layout_name_resolving_to_title_content_still_warns(tmp_path):
    # Unknown layout names resolve to the Title and Content layout (index
    # 1) — the warning gate must follow the RESOLVED layout, not the raw
    # spec string, or overflowing content escapes silently.
    spec = {
        "slide_size": "16:9",
        "slides": [
            {"layout": "title_conent",  # typo: resolves to title_content
             "title": "Dense",
             "bullets": [f"Bullet {i}: long operational detail sentence "
                         "that wraps in the placeholder frame and then "
                         "keeps going onto another rendered line"
                         for i in range(16)]},
        ],
    }
    result = run("pptx_create.py", write_spec(tmp_path, spec),
                 str(tmp_path / "typo.pptx"))
    assert result["ok"] is True
    assert result["overflow_warnings"], (
        "content on a slide that RESOLVES to title_content must warn even "
        "when the spec's layout name is unknown")
