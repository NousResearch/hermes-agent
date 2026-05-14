"""Safe Microsoft Office text extraction tool.

Narrow, read-only extractor for uploaded Office documents.  This deliberately
avoids shelling out to LibreOffice/pandoc so restricted gateway profiles can
read client briefs without enabling terminal or code_execution.
"""

from __future__ import annotations

import importlib.util
import json
from pathlib import Path
from typing import Any, Iterable

from tools.registry import registry, tool_error

SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".xlsx"}
LEGACY_EXTENSIONS = {".doc", ".ppt", ".xls"}
DEFAULT_MAX_CHARS = 80_000
DEFAULT_MAX_ROWS_PER_SHEET = 200
DEFAULT_MAX_COLS_PER_SHEET = 40


def _json(payload: dict[str, Any]) -> str:
    return json.dumps(payload, ensure_ascii=False)


def _clip(text: str, max_chars: int) -> tuple[str, bool]:
    if max_chars <= 0:
        max_chars = DEFAULT_MAX_CHARS
    if len(text) <= max_chars:
        return text, False
    marker = "\n\n[TRUNCATED: office_extract output exceeded max_chars]"
    return text[: max(0, max_chars - len(marker))].rstrip() + marker, True


def _error(message: str, *, error_code: str, file_path: str | None = None) -> str:
    payload: dict[str, Any] = {
        "success": False,
        "error": message,
        "error_code": error_code,
    }
    if file_path is not None:
        payload["file_path"] = file_path
    return _json(payload)


def _iter_nonempty(values: Iterable[Any]) -> list[str]:
    return [str(value).strip() for value in values if value is not None and str(value).strip()]


def _markdown_table(rows: list[list[str]]) -> list[str]:
    if not rows:
        return []
    width = max(len(row) for row in rows)
    normalized = [row + [""] * (width - len(row)) for row in rows]
    lines = [" | ".join(normalized[0])]
    if len(normalized) > 1:
        lines.append(" | ".join(["---"] * width))
        lines.extend(" | ".join(row) for row in normalized[1:])
    return lines


def _extract_docx(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        import docx
    except ImportError as exc:  # pragma: no cover - exercised by check_fn in normal use
        raise RuntimeError("python-docx is not installed") from exc

    document = docx.Document(str(path))
    lines: list[str] = []

    paragraphs = 0
    headings = 0
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        paragraphs += 1
        style_name = getattr(getattr(paragraph, "style", None), "name", "") or ""
        if style_name.lower().startswith("heading"):
            level = 1
            for token in style_name.split():
                if token.isdigit():
                    level = max(1, min(6, int(token)))
                    break
            lines.append(f"{'#' * level} {text}")
            headings += 1
        else:
            lines.append(text)

    table_count = 0
    for index, table in enumerate(document.tables, start=1):
        rows: list[list[str]] = []
        for row in table.rows:
            cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
            if any(cells):
                rows.append(cells)
        if not rows:
            continue
        table_count += 1
        lines.append(f"\nTable {index}:")
        lines.extend(_markdown_table(rows))

    metadata = {
        "paragraph_count": paragraphs,
        "heading_count": headings,
        "table_count": table_count,
    }
    return "\n\n".join(lines).strip(), metadata


def _extract_pptx(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        import pptx
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("python-pptx is not installed") from exc

    presentation = pptx.Presentation(str(path))
    lines: list[str] = []
    nonempty_slides = 0

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = "\n".join(_iter_nonempty(par.text for par in shape.text_frame.paragraphs))
            if text:
                slide_lines.append(text)

        # python-pptx exposes notes on demand; some decks have no notes part.
        try:
            notes_frame = slide.notes_slide.notes_text_frame
            notes = "\n".join(_iter_nonempty(par.text for par in notes_frame.paragraphs))
            if notes:
                slide_lines.append(f"Speaker notes:\n{notes}")
        except Exception:
            pass

        if slide_lines:
            nonempty_slides += 1
            lines.append(f"## Slide {slide_index}\n" + "\n\n".join(slide_lines))

    metadata = {
        "slide_count": len(presentation.slides),
        "nonempty_slide_count": nonempty_slides,
    }
    return "\n\n".join(lines).strip(), metadata


def _extract_xlsx(
    path: Path,
    *,
    max_rows_per_sheet: int = DEFAULT_MAX_ROWS_PER_SHEET,
    max_cols_per_sheet: int = DEFAULT_MAX_COLS_PER_SHEET,
) -> tuple[str, dict[str, Any]]:
    try:
        import openpyxl
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("openpyxl is not installed") from exc

    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    lines: list[str] = []
    sheet_meta: list[dict[str, Any]] = []

    for sheet in workbook.worksheets:
        rows: list[list[str]] = []
        for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            if row_index > max_rows_per_sheet:
                break
            cells = ["" if cell is None else str(cell).strip() for cell in row[:max_cols_per_sheet]]
            # Keep header-like rows and rows with at least one populated cell.
            if any(cells):
                rows.append(cells)

        if rows:
            lines.append(f"## Sheet: {sheet.title}")
            lines.extend(_markdown_table(rows))

        sheet_meta.append(
            {
                "name": sheet.title,
                "max_row": sheet.max_row,
                "max_column": sheet.max_column,
                "extracted_rows": len(rows),
                "truncated_rows": sheet.max_row > max_rows_per_sheet,
                "truncated_columns": sheet.max_column > max_cols_per_sheet,
            }
        )

    workbook.close()
    metadata = {
        "sheet_count": len(sheet_meta),
        "sheets": sheet_meta,
        "max_rows_per_sheet": max_rows_per_sheet,
        "max_cols_per_sheet": max_cols_per_sheet,
    }
    return "\n\n".join(lines).strip(), metadata


def office_extract(
    file_path: str,
    max_chars: int = DEFAULT_MAX_CHARS,
    max_rows_per_sheet: int = DEFAULT_MAX_ROWS_PER_SHEET,
    max_cols_per_sheet: int = DEFAULT_MAX_COLS_PER_SHEET,
    task_id: str | None = None,
) -> str:
    """Extract readable text/markdown from .docx, .pptx, or .xlsx files."""
    if not file_path or not str(file_path).strip():
        return _error("file_path is required", error_code="missing_file_path")

    path = Path(file_path).expanduser()
    if not path.exists():
        return _error(f"File not found: {file_path}", error_code="file_not_found", file_path=str(path))
    if not path.is_file():
        return _error(f"Path is not a file: {file_path}", error_code="not_a_file", file_path=str(path))

    suffix = path.suffix.lower()
    if suffix in LEGACY_EXTENSIONS:
        return _error(
            f"Legacy binary Office format '{suffix}' is not supported safely. "
            "Ask for .docx/.pptx/.xlsx or PDF instead.",
            error_code="unsupported_legacy_format",
            file_path=str(path),
        )
    if suffix not in SUPPORTED_EXTENSIONS:
        return _error(
            "Unsupported file type. Supported: .docx, .pptx, .xlsx. "
            "Legacy .doc/.ppt/.xls require conversion first.",
            error_code="unsupported_format",
            file_path=str(path),
        )

    try:
        if suffix == ".docx":
            markdown, metadata = _extract_docx(path)
        elif suffix == ".pptx":
            markdown, metadata = _extract_pptx(path)
        else:
            markdown, metadata = _extract_xlsx(
                path,
                max_rows_per_sheet=max_rows_per_sheet,
                max_cols_per_sheet=max_cols_per_sheet,
            )
    except Exception as exc:
        return tool_error(
            f"Failed to extract Office document '{path.name}': {exc}",
            success=False,
        )

    clipped, truncated = _clip(markdown, max_chars)
    return _json(
        {
            "success": True,
            "file_path": str(path),
            "format": suffix.lstrip("."),
            "markdown": clipped,
            "truncated": truncated,
            "metadata": metadata,
        }
    )


def check_office_extract_requirements() -> bool:
    return all(importlib.util.find_spec(name) is not None for name in ("docx", "pptx", "openpyxl"))


OFFICE_EXTRACT_SCHEMA = {
    "name": "office_extract",
    "description": (
        "Read Microsoft Office files from local paths and extract markdown/text. "
        "Supports .docx, .pptx, and .xlsx. Use this when a user uploads a Word, "
        "PowerPoint, or Excel file and you need its contents before answering. "
        "Does not execute macros and does not support legacy .doc/.ppt/.xls."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Absolute or relative path to a .docx, .pptx, or .xlsx file.",
            },
            "max_chars": {
                "type": "integer",
                "description": "Maximum characters of markdown to return. Default 80000.",
                "default": DEFAULT_MAX_CHARS,
            },
            "max_rows_per_sheet": {
                "type": "integer",
                "description": "For .xlsx only: max rows to extract per sheet. Default 200.",
                "default": DEFAULT_MAX_ROWS_PER_SHEET,
            },
            "max_cols_per_sheet": {
                "type": "integer",
                "description": "For .xlsx only: max columns to extract per sheet. Default 40.",
                "default": DEFAULT_MAX_COLS_PER_SHEET,
            },
        },
        "required": ["file_path"],
    },
}


registry.register(
    name="office_extract",
    toolset="office",
    schema=OFFICE_EXTRACT_SCHEMA,
    handler=lambda args, **kwargs: office_extract(
        file_path=args.get("file_path", ""),
        max_chars=args.get("max_chars", DEFAULT_MAX_CHARS),
        max_rows_per_sheet=args.get("max_rows_per_sheet", DEFAULT_MAX_ROWS_PER_SHEET),
        max_cols_per_sheet=args.get("max_cols_per_sheet", DEFAULT_MAX_COLS_PER_SHEET),
        task_id=kwargs.get("task_id"),
    ),
    check_fn=check_office_extract_requirements,
    description=OFFICE_EXTRACT_SCHEMA["description"],
    emoji="📄",
)
